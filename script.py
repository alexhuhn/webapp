from pptx import Presentation
import streamlit as st
from io import BytesIO

def create_ppt(x):
	prs = Presentation()
	title_slide_layout = prs.slide_layouts[0]
	slide = prs.slides.add_slide(title_slide_layout)
	title = slide.shapes.title
	subtitle = slide.placeholders[1]

	title.text = "Hello, World!"
	subtitle.text = "python-pptx was here!" + x
	
	return prs